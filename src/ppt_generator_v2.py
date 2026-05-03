#!/usr/bin/env python3
"""
Enhanced PPT Generator for μ-MuZero
====================================
Generates two versions:
  1. English version with formulas and figures
  2. Chinese version with formulas and figures
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Paths
BASE_DIR = os.path.expanduser("~/klaus_ai/mu-muzero")
FORMULA_DIR = os.path.join(BASE_DIR, "assets", "formulas")
FIGURE_DIR = os.path.join(BASE_DIR, "figures", "PhD")
MRES_FIG_DIR = os.path.join(BASE_DIR, "figures", "MRes")
SLIDES_DIR = os.path.join(BASE_DIR, "slides")
os.makedirs(SLIDES_DIR, exist_ok=True)

# Colors
DARK_BLUE = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT_BLUE = RGBColor(0x1E, 0x5E, 0xAA)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
ACCENT_ORANGE = RGBColor(0xE8, 0x5D, 0x04)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)


def create_presentation(lang="en"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_gradient_background(slide, prs, color1, color2=None):
    """Add a solid color background."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color1
    bg.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_title_slide(prs, title, subtitle, lang="en"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, prs, DARK_BLUE)
    
    # Decorative accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), Inches(0.15), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()
    
    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12), Inches(1))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)
    
    # Bottom info
    if lang == "en":
        info = "Hamlyn Centre for Robotic Surgery | Imperial College London | 2026"
    else:
        info = "哈姆林机器人手术中心 | 帝国理工学院 | 2026"
    box3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.4))
    tf3 = box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = info
    p3.font.size = Pt(12)
    p3.font.color.rgb = GRAY
    
    return slide


def add_section_header(prs, number, title, lang="en"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide, prs, ACCENT_BLUE)
    
    # Big number
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(3), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{number}" if number < 10 else str(number)
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.opacity = 0.3
    
    # Title
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(1))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(44)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    
    return slide


def add_content_slide(prs, title, bullets, formula_img=None, figure_img=None, 
                      figure_caption=None, lang="en", two_col=False, 
                      left_title=None, left_bullets=None, right_title=None, right_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Light background
    add_gradient_background(slide, prs, LIGHT_GRAY)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Title
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12), Inches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_left = Inches(0.4)
    content_width = Inches(12.5)
    content_top = Inches(1.2)
    content_height = Inches(5.8)
    
    # If figure on right, adjust layout
    if figure_img:
        content_width = Inches(6.5)
    
    if two_col:
        # Left column
        lbox = slide.shapes.add_textbox(content_left, content_top, Inches(5.8), content_height)
        ltf = lbox.text_frame
        ltf.word_wrap = True
        p = ltf.paragraphs[0]
        p.text = left_title or ""
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_after = Pt(10)
        for bullet in (left_bullets or []):
            p = ltf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(16)
            p.font.color.rgb = BLACK
            p.space_after = Pt(8)
        
        # Divider
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), content_top, Inches(0.02), Inches(5.0))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        line.line.fill.background()
        
        # Right column
        rbox = slide.shapes.add_textbox(Inches(6.6), content_top, Inches(5.8), content_height)
        rtf = rbox.text_frame
        rtf.word_wrap = True
        p = rtf.paragraphs[0]
        p.text = right_title or ""
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE
        p.space_after = Pt(10)
        for bullet in (right_bullets or []):
            p = rtf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(16)
            p.font.color.rgb = BLACK
            p.space_after = Pt(8)
    else:
        # Content bullets
        box2 = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        tf2 = box2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(18)
            p.font.color.rgb = BLACK
            p.space_after = Pt(10)
    
    # Add formula image
    if formula_img and os.path.exists(formula_img):
        slide.shapes.add_picture(formula_img, Inches(0.5), Inches(5.8), width=Inches(8))
    
    # Add figure image
    if figure_img and os.path.exists(figure_img):
        fig_left = Inches(7.0) if not two_col else Inches(9.0)
        fig_top = Inches(1.5) if not two_col else Inches(1.3)
        fig_width = Inches(5.5) if not two_col else Inches(3.8)
        pic = slide.shapes.add_picture(figure_img, fig_left, fig_top, width=fig_width)
        if figure_caption:
            cap_box = slide.shapes.add_textbox(fig_left, fig_top + pic.height + Inches(0.1), 
                                               fig_width, Inches(0.4))
            cap_tf = cap_box.text_frame
            cap_p = cap_tf.paragraphs[0]
            cap_p.text = figure_caption
            cap_p.font.size = Pt(10)
            cap_p.font.color.rgb = GRAY
            cap_p.font.italic = True
    
    return slide


def build_presentation(lang="en"):
    prs = create_presentation(lang)
    
    if lang == "en":
        T = {
            "title": "μ-MuZero: Model-Based RL for Autonomous Micro-Robotic Surgery",
            "subtitle": "A Safety-Constrained, Stochastic Variant of MuZero for Optical Micromanipulation",
            "sec1": "Motivation & Background",
            "sec2": "Core Algorithm",
            "sec3": "System Architecture",
            "sec4": "Experimental Validation",
            "sec5": "Future & Impact",
            "slide2_title": "Why We Need μ-MuZero",
            "slide2_bullets": [
                "Optical microrobots can perform cell-level surgery, but are still manually operated",
                "Human reaction time (~250 ms) is too slow for micro-scale fluid dynamics",
                "Standard RL algorithms fail at micro-scale due to three unique challenges:",
                "   1. Brownian motion makes dynamics inherently stochastic",
                "   2. Laser power has hard safety bounds — cannot explore by trial-and-error",
                "   3. Action space is combinatorial: 6 traps × 7 primitives = 117,649 actions",
                "MuZero works for board games, but not for micro-robots without adaptation",
            ],
            "slide3_title": "The Core Insight: What Changes at Micro-Scale?",
            "slide3_bullets": [
                "At macro-scale: inertia dominates, dynamics are deterministic, safety is soft",
                "At micro-scale (low Reynolds number): viscosity dominates, Brownian motion is irreducible",
                "Optical force fields are highly non-linear — no closed-form model exists",
                "Key realization: agent must reason about UNCERTAINTY, not just expectation",
                "μ-MuZero learns a stochastic latent dynamics model + plans with safety-constrained MCTS",
            ],
            "slide4_title": "μ-MuZero Architecture",
            "slide4_left": "Neural Networks",
            "slide4_left_bullets": [
                "Representation h_θ: belief → latent state (256-dim)",
                "Dynamics g_θ: (state, action) → next state distribution",
                "   Outputs μ and σ for Gaussian stochastic transitions",
                "Prediction f_θ: state → policy logits + value estimate",
                "All trained end-to-end via self-play in digital twin",
            ],
            "slide4_right": "Safety-Constrained MCTS",
            "slide4_right_bullets": [
                "Standard UCB augmented with uncertainty penalty",
                "Hard pruning: any action violating P > P_max is removed",
                "Soft penalty: phototoxicity budget tracked per episode",
                "Stochastic rollouts: sample K=8 trajectories from model",
                "Value averaged across samples for risk-sensitive planning",
            ],
            "slide5_title": "Innovation 1: Stochastic Latent Dynamics",
            "slide5_bullets": [
                "Standard MuZero: deterministic transitions s' = g(s, a)",
                "μ-MuZero: stochastic transitions s' ~ N(μ_g, σ_g)",
                "Brownian displacement in 50 ms ≈ 100 nm — comparable to trap motion",
                "During planning, sample K trajectories and average values",
                "Result: policy robust to irreducible noise",
            ],
            "slide5_formula": "stochastic_dynamics",
            "slide6_title": "Innovation 2: Safety-Constrained Tree Search",
            "slide6_bullets": [
                "In surgery, 'explore then recover' is unacceptable",
                "Level 1 (Hard): MCTS prunes nodes where P > P_max",
                "Level 2 (Soft): cumulative phototoxicity budget penalised",
                "Level 3 (Epistemic): high uncertainty reduces exploration",
                "Constraint violation: 0.3% (μ-MuZero) vs 18% (unconstrained)",
            ],
            "slide6_formula": "ucb",
            "slide7_title": "Innovation 3: Factored Action Representation",
            "slide7_bullets": [
                "Naive space: 7^6 = 117,649 joint actions — intractable",
                "μ-MuZero factorises: select WHICH trap, then HOW to move",
                "Branching factor: 7^K → K + 6 (e.g., 117,649 → 12)",
                "Biologically plausible: humans attend to one trap at a time",
                "Enables real-time planning at 20 Hz on standard GPU",
            ],
            "slide7_formula": "factored_action",
            "slide8_title": "Curriculum Learning for Multi-Robot Coordination",
            "slide8_bullets": [
                "Stage 1 (0–50K): Single-robot transport in empty workspace",
                "Stage 2 (50K–250K): Add static obstacles",
                "Stage 3 (250K–1M): Two-robot cooperative pushing",
                "Stage 4 (1M–2M): Three-robot micro-assembly",
                "Without curriculum: fails to converge after 5M steps",
            ],
            "slide9_title": "Application Scenarios in Micro-Surgery",
            "slide9_left": "Single-Robot Tasks",
            "slide9_left_bullets": [
                "Cell transport: move target to micro-channel entrance",
                "Obstacle avoidance: navigate through cluttered tissue",
                "Cell sorting: separate target cells from debris",
                "Success rate: 91.3% (vs 67.2% scripted, 54.8% PPO)",
            ],
            "slide9_right": "Multi-Robot Tasks",
            "slide9_right_bullets": [
                "Cooperative transport: 3 robots push untrappable cell",
                "Micro-assembly: position component into slot",
                "Coordinated injection: multi-point drug delivery",
                "Success rate: 78.4% (scripted baseline: 0%)",
            ],
            "slide10_title": "Physics-Informed Digital Twin",
            "slide10_bullets": [
                "Learned optical force model (NN trained on FDTD/T-matrix)",
                "Hydrodynamic interactions via RPY tensor",
                "Brownian dynamics integrated with Euler-Maruyama",
                "Differentiable imaging simulator for sim-to-real transfer",
            ],
            "slide10_formula": "langevin",
            "slide11_title": "Performance Summary vs Baselines",
            "slide11_bullets": [
                "Task 1 (Transport): μ-MuZero 91.3% | Scripted 67.2% | PPO 54.8% | Oracle 94.1%",
                "Task 2 (Obstacle): μ-MuZero 89.0% | Scripted 58.0% | Manual 62.0%",
                "Task 3 (Cooperative): μ-MuZero 78.4% | Scripted 0% | Manual 34.0%",
                "Task 4 (Assembly): μ-MuZero 68.0% | Scripted 0% | Manual 8.0%",
                "Phototoxicity budget: μ-MuZero 45% | Manual 68% | Scripted 82%",
                "Emergent behaviour: discovers 'sling-shot' manoeuvre",
            ],
            "slide12_title": "Sim-to-Real Transfer Pathway",
            "slide12_bullets": [
                "Step 1: Domain randomisation — randomise viscosity, temperature, noise",
                "Step 2: System identification — calibrate on 100 real trajectories",
                "Step 3: Online adaptation — MPC compensates residual model error",
                "Expected transfer loss: < 15% success rate drop",
                "Physical validation: next critical milestone at Hamlyn Centre",
            ],
            "slide13_title": "Future Directions & Open Problems",
            "slide13_bullets": [
                "Real-world validation on physical optical tweezer hardware",
                "Meta-learning for rapid adaptation to new microrobot geometries",
                "Hierarchical planning: strategic (seconds) + tactical (ms) levels",
                "Interpretability: explain WHY the agent chose a movement",
                "Integration with fibre-optic endoscopes for in vivo deployment",
                "Regulatory pathway: certifying learned policies for human surgery",
            ],
            "slide14_title": "References & Resources",
            "slide14_bullets": [
                "Schrittwieser et al. (2020). Mastering Atari, Go, chess and shogi by planning with a learned model. Nature.",
                "Ren et al. (2022). Machine learning-based real-time localization and automatic trapping of multiple microrobots. MARSS.",
                "Zhang et al. (2020). Distributed force control for microrobot manipulation via planar multi-spot optical tweezer. AOM.",
                "Grammatikopoulou et al. (2017). Depth estimation of optically transparent laser-driven microrobots. IROS.",
                "Code: github.com/yunxiao-ren/mu-muzero (upon publication)",
                "Contact: yunxiao.ren@imperial.ac.uk",
            ],
            "fig_arch_caption": "EOMF Three-Layer Architecture",
            "fig_train_caption": "Training Curves: μ-MuZero vs PPO",
            "fig_ablation_caption": "Ablation Study Results",
            "fig_mcts_caption": "Safety-Constrained MCTS Schematic",
        }
    else:
        T = {
            "title": "μ-MuZero：面向自主微机器人手术的模型强化学习",
            "subtitle": "面向光学微操作的 MuZero 安全约束随机变体",
            "sec1": "研究动机与背景",
            "sec2": "核心算法",
            "sec3": "系统架构",
            "sec4": "实验验证",
            "sec5": "未来展望",
            "slide2_title": "为什么需要 μ-MuZero？",
            "slide2_bullets": [
                "光学微机器人可执行细胞级手术，但目前仍需手动操作",
                "人类反应时间 (~250 ms) 对微尺度流体动力学来说太慢",
                "标准强化学习在微尺度面临三大独特挑战：",
                "   1. 布朗运动使动力学本质随机",
                "   2. 激光功率有硬安全边界 — 不能通过试错探索",
                "   3. 动作空间组合爆炸：6个光阱 × 7种原语 = 117,649种动作",
                "MuZero 适用于棋盘游戏，但需改造才能用于微机器人",
            ],
            "slide3_title": "核心洞察：微尺度发生了什么本质变化？",
            "slide3_bullets": [
                "宏观尺度：惯性主导、确定性动力学、安全是软约束",
                "微尺度（低雷诺数）：粘性主导、布朗运动不可忽略",
                "光力场高度非线性 — 不存在闭式动力学模型",
                "关键认识：智能体必须推理不确定性，而非仅期望值",
                "μ-MuZero 学习随机潜态动力学模型 + 安全约束 MCTS 规划",
            ],
            "slide4_title": "μ-MuZero 算法架构",
            "slide4_left": "神经网络",
            "slide4_left_bullets": [
                "表征网络 h_θ：信念 → 潜态 (256维)",
                "动力学网络 g_θ：(状态, 动作) → 下一状态分布",
                "   输出 μ 和 σ 用于高斯随机转移",
                "预测网络 f_θ：状态 → 策略逻辑值 + 价值估计",
                "全部通过数字孪生中的自博弈端到端训练",
            ],
            "slide4_right": "安全约束 MCTS",
            "slide4_right_bullets": [
                "标准 UCB 公式增加不确定性惩罚项",
                "硬剪枝：违反 P > P_max 的动作直接移除",
                "软惩罚：每回合跟踪光毒性预算",
                "随机推出：从学习模型采样 K=8 条轨迹",
                "跨样本平均价值实现风险敏感规划",
            ],
            "slide5_title": "创新点 1：随机潜态动力学",
            "slide5_bullets": [
                "标准 MuZero：确定性转移 s' = g(s, a)",
                "μ-MuZero：随机转移 s' ~ N(μ_g, σ_g)",
                "50 ms 内布朗运动位移 ≈ 100 nm — 与光阱运动相当",
                "规划时采样 K 条轨迹并平均价值",
                "结果：策略对不可约噪声具有鲁棒性",
            ],
            "slide5_formula": "stochastic_dynamics",
            "slide6_title": "创新点 2：安全约束树搜索",
            "slide6_bullets": [
                "手术中'探索后恢复'不可接受 — 一次过功率事件即可杀死细胞",
                "第一层（硬约束）：MCTS 剪枝 P > P_max 的节点",
                "第二层（软约束）：累积光毒性预算加入奖励惩罚",
                "第三层（认知层）：高感知不确定性降低探索奖励",
                "约束违反率：0.3% (μ-MuZero) vs 18% (无约束基线)",
            ],
            "slide6_formula": "ucb",
            "slide7_title": "创新点 3：分解动作表示",
            "slide7_bullets": [
                "朴素空间：7^6 = 117,649 种联合动作 — 无法处理",
                "μ-MuZero 分解：先选哪个光阱，再选如何移动",
                "分支因子：7^K → K + 6（如 117,649 → 12）",
                "生物学合理：人类也一次只关注一个光阱",
                "在标准 GPU 上实现 20 Hz 实时规划",
            ],
            "slide7_formula": "factored_action",
            "slide8_title": "多机器人协调的课程学习",
            "slide8_bullets": [
                "第一阶段 (0–50K 步)：空工作区单机器人运输",
                "第二阶段 (50K–250K)：加入静态障碍物",
                "第三阶段 (250K–1M)：双机器人协作推运",
                "第四阶段 (1M–2M)：三机器人微装配",
                "无课程学习：即使 5M 步后仍无法收敛",
            ],
            "slide9_title": "微手术应用场景",
            "slide9_left": "单机器人任务",
            "slide9_left_bullets": [
                "细胞运输：将目标细胞移至微通道入口",
                "避障导航：在杂乱组织中穿行",
                "细胞分选：从碎片中分离目标细胞",
                "成功率：91.3%（对比脚本 67.2%，PPO 54.8%）",
            ],
            "slide9_right": "多机器人任务",
            "slide9_right_bullets": [
                "协作运输：3个机器人推无法直接捕获的大细胞",
                "微装配：将组件精确插入卡槽",
                "协调注射：多点同时给药",
                "成功率：78.4%（脚本基线：0%）",
            ],
            "slide10_title": "物理信息数字孪生",
            "slide10_bullets": [
                "学习光力模型（在 FDTD/T-matrix 数据上训练的神经网络）",
                "通过 RPY 张量建模流体动力学相互作用",
                "欧拉-丸山法积分布朗动力学",
                "可微分成像模拟器实现仿真到现实迁移",
            ],
            "slide10_formula": "langevin",
            "slide11_title": "与基线方法的性能对比",
            "slide11_bullets": [
                "任务1 (运输)：μ-MuZero 91.3% | 脚本 67.2% | PPO 54.8% | 预言 94.1%",
                "任务2 (避障)：μ-MuZero 89.0% | 脚本 58.0% | 手动 62.0%",
                "任务3 (协作)：μ-MuZero 78.4% | 脚本 0% | 手动 34.0%",
                "任务4 (装配)：μ-MuZero 68.0% | 脚本 0% | 手动 8.0%",
                "光毒性预算：μ-MuZero 45% | 手动 68% | 脚本 82%",
                "涌现行为：发现脚本控制器中不存在的'弹弓'策略",
            ],
            "slide12_title": "仿真到现实的迁移路径",
            "slide12_bullets": [
                "第一步：域随机化 — 训练时随机化粘度、温度、噪声",
                "第二步：系统辨识 — 用 100 条真实轨迹校准孪生参数",
                "第三步：在线适应 — MPC 补偿残余模型误差",
                "预期迁移损失：成功率下降 < 15%",
                "物理验证：哈姆林中心的下一个关键里程碑",
            ],
            "slide13_title": "未来方向与开放问题",
            "slide13_bullets": [
                "在物理光学镊硬件上进行真实世界验证",
                "元学习实现对新微机器人几何的快速适应",
                "层级规划：战略层（秒级）+ 战术层（毫秒级）",
                "可解释性：解释智能体为何选择某个移动",
                "与光纤内窥镜集成实现体内部署",
                "监管路径：如何为人体手术认证学习策略",
            ],
            "slide14_title": "参考文献与资源",
            "slide14_bullets": [
                "Schrittwieser et al. (2020). Mastering Atari, Go, chess and shogi by planning with a learned model. Nature.",
                "Ren et al. (2022). Machine learning-based real-time localization and automatic trapping of multiple microrobots. MARSS.",
                "Zhang et al. (2020). Distributed force control for microrobot manipulation via planar multi-spot optical tweezer. AOM.",
                "Grammatikopoulou et al. (2017). Depth estimation of optically transparent laser-driven microrobots. IROS.",
                "代码：github.com/yunxiao-ren/mu-muzero（发表后公开）",
                "联系：yunxiao.ren@imperial.ac.uk",
            ],
            "fig_arch_caption": "EOMF 三层架构",
            "fig_train_caption": "训练曲线：μ-MuZero vs PPO",
            "fig_ablation_caption": "消融实验结果",
            "fig_mcts_caption": "安全约束 MCTS 示意图",
        }
    
    # Build slides
    add_title_slide(prs, T["title"], T["subtitle"], lang)
    
    # Section 1: Motivation
    add_section_header(prs, 1, T["sec1"], lang)
    add_content_slide(prs, T["slide2_title"], T["slide2_bullets"], 
                     figure_img=os.path.join(FIGURE_DIR, "fig_eomf_arch.pdf"),
                     figure_caption=T["fig_arch_caption"], lang=lang)
    add_content_slide(prs, T["slide3_title"], T["slide3_bullets"],
                     formula_img=os.path.join(FORMULA_DIR, "pomdp.png"), lang=lang)
    
    # Section 2: Core Algorithm
    add_section_header(prs, 2, T["sec2"], lang)
    add_content_slide(prs, T["slide4_title"], [], two_col=True,
                     left_title=T["slide4_left"], left_bullets=T["slide4_left_bullets"],
                     right_title=T["slide4_right"], right_bullets=T["slide4_right_bullets"], lang=lang)
    add_content_slide(prs, T["slide5_title"], T["slide5_bullets"],
                     formula_img=os.path.join(FORMULA_DIR, "stochastic_dynamics.png"),
                     figure_img=os.path.join(FIGURE_DIR, "fig_training_curve.pdf"),
                     figure_caption=T["fig_train_caption"], lang=lang)
    add_content_slide(prs, T["slide6_title"], T["slide6_bullets"],
                     formula_img=os.path.join(FORMULA_DIR, "ucb.png"),
                     figure_img=os.path.join(FIGURE_DIR, "fig_mcts_tree.pdf"),
                     figure_caption=T["fig_mcts_caption"], lang=lang)
    add_content_slide(prs, T["slide7_title"], T["slide7_bullets"],
                     formula_img=os.path.join(FORMULA_DIR, "factored_action.png"), lang=lang)
    add_content_slide(prs, T["slide8_title"], T["slide8_bullets"],
                     figure_img=os.path.join(FIGURE_DIR, "fig_curriculum.pdf"),
                     figure_caption="Curriculum Learning Stages", lang=lang)
    
    # Section 3: Architecture
    add_section_header(prs, 3, T["sec3"], lang)
    add_content_slide(prs, T["slide10_title"], T["slide10_bullets"],
                     formula_img=os.path.join(FORMULA_DIR, "langevin.png"),
                     figure_img=os.path.join(FIGURE_DIR, "fig_digital_twin.pdf"),
                     figure_caption="Physics-Informed Digital Twin", lang=lang)
    add_content_slide(prs, T["slide9_title"], [], two_col=True,
                     left_title=T["slide9_left"], left_bullets=T["slide9_left_bullets"],
                     right_title=T["slide9_right"], right_bullets=T["slide9_right_bullets"], lang=lang)
    
    # Section 4: Validation
    add_section_header(prs, 4, T["sec4"], lang)
    add_content_slide(prs, T["slide11_title"], T["slide11_bullets"],
                     figure_img=os.path.join(FIGURE_DIR, "fig_validation_bar.pdf"),
                     figure_caption="Task Success Rates", lang=lang)
    add_content_slide(prs, T["slide11_title"] + " (2)", [],
                     figure_img=os.path.join(FIGURE_DIR, "fig_ablation.pdf"),
                     figure_caption=T["fig_ablation_caption"], lang=lang)
    add_content_slide(prs, T["slide12_title"], T["slide12_bullets"],
                     figure_img=os.path.join(FIGURE_DIR, "fig_sim2real_gap.pdf"),
                     figure_caption="Sim-to-Real Gap", lang=lang)
    
    # Section 5: Future
    add_section_header(prs, 5, T["sec5"], lang)
    add_content_slide(prs, T["slide13_title"], T["slide13_bullets"], lang=lang)
    add_content_slide(prs, T["slide14_title"], T["slide14_bullets"], lang=lang)
    
    # Save
    suffix = "zh" if lang == "zh" else "en"
    path = os.path.join(SLIDES_DIR, f"mu_muzero_algorithm_{suffix}.pptx")
    prs.save(path)
    print(f"Saved: {path} ({len(prs.slides)} slides)")
    return path


if __name__ == "__main__":
    build_presentation("en")
    build_presentation("zh")
    print("\nBoth versions generated successfully!")
