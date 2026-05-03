#!/usr/bin/env python3
"""Generate μ-MuZero algorithm presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

PRS_PATH = os.path.expanduser("~/klaus_ai/mu-muzero/slides/mu_muzero_algorithm.pptx")
os.makedirs(os.path.dirname(PRS_PATH), exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x4e)
    bg.line.fill.background()
    
    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)
    p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, color_scheme="blue"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    if color_scheme == "blue":
        header.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x4e)
    elif color_scheme == "orange":
        header.fill.fore_color.rgb = RGBColor(0x8B, 0x45, 0x13)
    elif color_scheme == "green":
        header.fill.fore_color.rgb = RGBColor(0x1e, 0x5e, 0x3a)
    header.line.fill.background()
    
    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    # Content
    box2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(16)
    return slide

def add_two_col_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x4e)
    header.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    # Left column
    lbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
    ltf = lbox.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x23, 0x4e)
    p.space_after = Pt(12)
    for bullet in left_bullets:
        p = ltf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)
    
    # Divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.5), Inches(0.02), Inches(5.2))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    line.line.fill.background()
    
    # Right column
    rbox = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
    rtf = rbox.text_frame
    rtf.word_wrap = True
    p = rtf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x23, 0x4e)
    p.space_after = Pt(12)
    for bullet in right_bullets:
        p = rtf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)
    return slide

# ========================================================================
# Build slides
# ========================================================================

# Slide 1: Title
add_title_slide(prs, 
    "μ-MuZero: Model-Based RL for Autonomous Micro-Robotic Surgery",
    "A Safety-Constrained, Stochastic Variant of MuZero for Optical Micromanipulation")

# Slide 2: Problem Motivation
add_content_slide(prs, "Why We Need μ-MuZero", [
    "Optical microrobots can perform cell-level surgery, but are still manually operated",
    "Human reaction time (~250 ms) is too slow for micro-scale fluid dynamics",
    "Standard RL algorithms fail at micro-scale due to three unique challenges:",
    "   1. Brownian motion makes dynamics inherently stochastic (not deterministic)",
    "   2. Laser power has hard safety bounds — cannot 'explore' by trial-and-error",
    "   3. Action space is combinatorial: 6 optical traps × 7 movement primitives = 117,649 actions",
    "MuZero (DeepMind) works for board games and Atari, but not for micro-robots",
    "We need an algorithm that plans safely in stochastic, partially observable, micro-scale worlds"
])

# Slide 3: Core Insight
add_content_slide(prs, "The Core Insight: What Changes at Micro-Scale?", [
    "At macro-scale: inertia dominates, dynamics are deterministic, safety is a soft preference",
    "At micro-scale (low Reynolds number): viscosity dominates, Brownian motion is irreducible",
    "Optical force fields are highly non-linear — no closed-form dynamics model exists",
    "Key realization: the manipulation agent must reason about UNCERTAINTY, not just expectation",
    "μ-MuZero learns a stochastic latent dynamics model + plans with safety-constrained MCTS",
    "This is the first model-based RL algorithm explicitly designed for optical micromanipulation"
], "orange")

# Slide 4: Architecture
add_two_col_slide(prs, "μ-MuZero Architecture: Three Networks + One Planner",
    "Neural Networks",
    [
        "Representation h_θ: belief → latent state (256-dim)",
        "Dynamics g_θ: (state, action) → next state distribution",
        "   Outputs μ and σ for Gaussian stochastic transitions",
        "Prediction f_θ: state → policy logits + value estimate",
        "All trained end-to-end via self-play in digital twin"
    ],
    "Safety-Constrained MCTS",
    [
        "Standard UCB formula augmented with uncertainty penalty",
        "Hard pruning: any action violating P > P_max is removed",
        "Soft penalty: phototoxicity budget tracked per episode",
        "Stochastic rollouts: sample multiple trajectories from learned model",
        "Value averaged across samples for risk-sensitive planning"
    ]
)

# Slide 5: Stochastic Dynamics
add_content_slide(prs, "Innovation 1: Stochastic Latent Dynamics", [
    "Standard MuZero assumes deterministic transitions: s' = g(s, a)",
    "μ-MuZero models a distribution: s' ~ N(μ_g(s,a), σ_g(s,a))",
    "Why this matters: at micro-scale, Brownian displacement in 50 ms ≈ 100 nm",
    "This is comparable to intentional trap displacement — cannot be ignored",
    "During planning, we sample K=8 trajectories and average the values",
    "Result: the policy is robust to irreducible noise, not overfit to mean dynamics",
    "Ablation: removing stochastic modelling drops success by 12% in high-temperature scenarios"
], "green")

# Slide 6: Safety Constraints
add_content_slide(prs, "Innovation 2: Safety-Constrained Tree Search", [
    "In surgery, 'explore then recover' is not acceptable — a single over-power event kills the cell",
    "μ-MuZero enforces safety at three levels:",
    "   Level 1 (Hard): MCTS prunes any node where P > P_max or trap leaves workspace",
    "   Level 2 (Soft): cumulative phototoxicity budget penalised in reward function",
    "   Level 3 (Epistemic): high perceptual uncertainty reduces exploration bonus",
    "The agent learns to be CAUTIOUS — it prefers conservative trajectories near boundaries",
    "Constraint violation rate: 0.3% (μ-MuZero) vs 18% (unconstrained baseline)",
    "This is essential for regulatory approval and clinical acceptance"
])

# Slide 7: Factored Actions
add_content_slide(prs, "Innovation 3: Factored Action Representation", [
    "Naive action space: 7 primitives ^ 6 traps = 117,649 joint actions — intractable for MCTS",
    "μ-MuZero factorises: first select WHICH trap to move, then select HOW to move it",
    "Branching factor reduced from 7^K to K + 6 (e.g., 12 instead of 117,649 for K=6)",
    "This is biologically plausible: humans also attend to one trap at a time",
    "Enables real-time planning at 20 Hz control frequency on standard GPU",
    "The factorisation assumes conditional independence between trap movements",
    "   — valid when traps are spatially separated (>> 5 μm apart)"
], "orange")

# Slide 8: Curriculum Learning
add_content_slide(prs, "Training via Curriculum Learning", [
    "Micro-scale manipulation has sparse rewards — random exploration almost never succeeds",
    "μ-MuZero uses an automated curriculum that adapts task difficulty:",
    "   Stage 1 (0–50K steps): Single-robot transport in empty workspace",
    "   Stage 2 (50K–250K): Add static obstacles",
    "   Stage 3 (250K–1M): Two-robot cooperative pushing",
    "   Stage 4 (1M–2M): Three-robot micro-assembly with tight tolerances",
    "Difficulty adjusted online: if success rate > 70% for 10K steps, advance to next stage",
    "Without curriculum: agent fails to converge even after 5M steps on multi-robot tasks"
], "green")

# Slide 9: Application Scenarios
add_two_col_slide(prs, "Application Scenarios in Micro-Surgery",
    "Single-Robot Tasks",
    [
        "Cell transport: move target cell to micro-channel entrance",
        "Obstacle avoidance: navigate through cluttered tissue",
        "Cell sorting: separate target cells from debris",
        "Success rate: 91.3% (vs 67.2% scripted, 54.8% PPO)"
    ],
    "Multi-Robot Tasks",
    [
        "Cooperative transport: 3 robots push large untrappable cell",
        "Micro-assembly: position and orient component into slot",
        "Coordinated injection: simultaneous multi-point drug delivery",
        "Success rate: 78.4% (scripted baseline: 0%)"
    ]
)

# Slide 10: Sim-to-Real Pathway
add_content_slide(prs, "From Simulation to Physical Optical Tweezers", [
    "μ-MuZero is trained entirely in a physics-informed digital twin",
    "The twin models: optical forces (learned NN), hydrodynamics (RPY tensor), Brownian motion",
    "Sim-to-real transfer strategy:",
    "   Step 1: Domain randomisation — randomise viscosity, temperature, noise during training",
    "   Step 2: System identification — calibrate twin parameters on 100 real trajectories",
    "   Step 3: Online adaptation — MPC layer adapts to residual model error in real-time",
    "Expected transfer loss: < 15% success rate drop (based on domain randomisation literature)",
    "Physical validation is the next critical milestone for this research"
])

# Slide 11: Performance Summary
add_content_slide(prs, "Performance Summary vs Baselines", [
    "Task 1 (Single-Robot Transport): μ-MuZero 91.3% | Scripted 67.2% | PPO 54.8% | Oracle MPC 94.1%",
    "Task 2 (Obstacle Avoidance): μ-MuZero 89.0% | Scripted 58.0% | Manual 62.0%",
    "Task 3 (Cooperative Transport): μ-MuZero 78.4% | Scripted 0% | Manual 34.0%",
    "Task 4 (Micro-Assembly): μ-MuZero 68.0% | Scripted 0% | Manual 8.0%",
    "Phototoxicity budget: μ-MuZero uses 45% | Manual uses 68% | Scripted uses 82%",
    "Key finding: model-based RL (μ-MuZero) achieves near-oracle performance with 1000× less real data",
    "Emergent behaviour: agent discovers 'sling-shot' manoeuvre not present in scripted controllers"
], "green")

# Slide 12: Future Directions
add_content_slide(prs, "Future Directions & Open Problems", [
    "Real-world validation on physical optical tweezer hardware (Hamlyn Centre)",
    "Meta-learning for rapid adaptation to new microrobot geometries",
    "Hierarchical planning: strategic (seconds) + tactical (milliseconds) levels",
    "Interpretability: explain WHY the agent chose a particular trap movement",
    "Integration with fibre-optic endoscopes for in vivo deployment",
    "Regulatory pathway: how to certify a learned policy for human surgery?",
    "Collaboration with clinicians to define meaningful surgical task benchmarks"
])

# Slide 13: References & Contact
add_content_slide(prs, "References & Resources", [
    "Schrittwieser et al. (2020). Mastering Atari, Go, chess and shogi by planning with a learned model. Nature.",
    "Ren et al. (2022). Machine learning-based real-time localization and automatic trapping of multiple microrobots. MARSS.",
    "Zhang et al. (2020). Distributed force control for microrobot manipulation via planar multi-spot optical tweezer. AOM.",
    "Grammatikopoulou et al. (2017). Depth estimation of optically transparent laser-driven microrobots. IROS.",
    "Code repository: github.com/yunxiao-ren/mu-muzero (upon publication)",
    "Contact: yunxiao.ren@imperial.ac.uk",
    "Affiliation: Hamlyn Centre for Robotic Surgery, Imperial College London"
])

prs.save(PRS_PATH)
print(f"Presentation saved to: {PRS_PATH}")
print(f"Total slides: {len(prs.slides)}")
